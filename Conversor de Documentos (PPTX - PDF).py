import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter.ttk import Combobox
from pptx import Presentation
from PyPDF2 import PdfReader, PdfWriter
import os


def convert_file():
    input_format = input_format_combobox.get()
    output_format = output_format_combobox.get()
    file_path = file_path_label.cget("text")
    
    if input_format == output_format:
        messagebox.showwarning("Atenção", "Os formatos de entrada e saída são os mesmos!")
        return
    
    if not file_path or file_path == "Nenhum arquivo selecionado":
        messagebox.showerror("Erro", "Por favor, selecione um arquivo para converter.")
        return
    
    try:
        if input_format == "PDF" and output_format == "PPTX":
            pdf_to_pptx(file_path)
        elif input_format == "PPTX" and output_format == "PDF":
            pptx_to_pdf(file_path)
        else:
            messagebox.showerror("Erro", "Conversão não suportada.")
            return
        messagebox.showinfo("Sucesso", "Arquivo convertido com sucesso!")
    except Exception as e:
        messagebox.showerror("Erro", f"Erro durante a conversão: {e}")


def select_file():
    input_format = input_format_combobox.get()
    file_types = [("Todos os arquivos", "*.*")]
    if input_format == "PDF":
        file_types = [("Arquivos PDF", "*.pdf")]
    elif input_format == "PPTX":
        file_types = [("Arquivos PPTX", "*.pptx")]
    
    file_path = filedialog.askopenfilename(filetypes=file_types)
    if file_path:
        file_path_label.config(text=file_path)


def pdf_to_pptx(file_path):
    # Exemplo básico de conversão - Cada página do PDF será um slide
    reader = PdfReader(file_path)
    presentation = Presentation()

    for page in reader.pages:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        text = page.extract_text()
        textbox = slide.shapes.add_textbox(left=0, top=0, width=presentation.slide_width, height=presentation.slide_height)
        textbox.text = text

    output_path = os.path.splitext(file_path)[0] + "_converted.pptx"
    presentation.save(output_path)


def pptx_to_pdf(file_path):
    # Exemplo básico que cria um PDF simples com o texto de cada slide
    presentation = Presentation(file_path)
    writer = PdfWriter()

    for slide in presentation.slides:
        text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
        
        writer.add_blank_page()
        writer.pages[-1].mediaBox.upperRight = (300, 400)
        writer.pages[-1].insert_text(text)
    
    output_path = os.path.splitext(file_path)[0] + "_converted.pdf"
    with open(output_path, "wb") as f:
        writer.write(f)


# Configurando a interface
root = tk.Tk()
root.title("Conversor de Arquivos")
root.geometry("580x300")
root.configure(bg="#e0f7fa")

# Estilo do título
title_label = tk.Label(root, text="Conversor de Arquivos", bg="#00796b", fg="white", font=("Arial", 18, "bold"), pady=10)
title_label.pack(fill="x")

# Seleção de formatos
frame_options = tk.Frame(root, bg="#e0f7fa")
frame_options.pack(pady=10)

tk.Label(frame_options, text="Formato de entrada:", bg="#e0f7fa", font=("Arial", 12)).grid(row=0, column=0, padx=5, pady=5)
input_format_combobox = Combobox(frame_options, values=["PDF", "PPTX"], state="readonly", font=("Arial", 10))
input_format_combobox.grid(row=0, column=1, padx=5, pady=5)
input_format_combobox.set("PDF")

tk.Label(frame_options, text="Formato de saída:", bg="#e0f7fa", font=("Arial", 12)).grid(row=1, column=0, padx=5, pady=5)
output_format_combobox = Combobox(frame_options, values=["PDF", "PPTX"], state="readonly", font=("Arial", 10))
output_format_combobox.grid(row=1, column=1, padx=5, pady=5)
output_format_combobox.set("PPTX")

# Seleção de arquivo
file_selection_frame = tk.Frame(root, bg="#e0f7fa")
file_selection_frame.pack(pady=10)

file_path_label = tk.Label(file_selection_frame, text="Nenhum arquivo selecionado", bg="#b2dfdb", font=("Arial", 10), anchor="w", width=50)
file_path_label.pack(side="left", padx=5)

select_file_button = tk.Button(file_selection_frame, text="Escolher Arquivo", command=select_file, bg="#004d40", fg="white", font=("Arial", 10, "bold"))
select_file_button.pack(side="right", padx=5)

# Botão de conversão
convert_button = tk.Button(root, text="Converter", command=convert_file, bg="#00796b", fg="white", font=("Arial", 12, "bold"), pady=10)
convert_button.pack(pady=20)

root.mainloop()
